"""Conversion routes — Gotenberg for documents, Ghostscript for PDF compress."""

import asyncio
import base64
import io
import json
import mimetypes
import os
import posixpath
import re
import subprocess
import tempfile
import time
import zipfile
from collections import defaultdict
from xml.etree import ElementTree as ET

import httpx
from data.db import async_engine
from data.models import User
from data.security import current_user_for_convert, require_admin
from fastapi import APIRouter, Depends, File, Form, UploadFile
from fastapi.responses import Response
from log import get_logger, request_id_var
from sqlalchemy import text
from validation import (
    MAX_FILE_SIZE,
    ValidationError,
    validate_compress_quality,
    validate_upload,
)

logger = get_logger("converter")

router = APIRouter(prefix="/api/v1")

GOTENBERG_URL = os.getenv("GOTENBERG_URL", "http://gotenberg:3000")
REQUEST_TIMEOUT = 60.0
# /health's DB connectivity check (P2 §20) — bounds the same way the
# Gotenberg check's httpx timeout does. A module-level constant (not an
# inline literal) so a test can monkeypatch it down instead of actually
# waiting out a multi-second timeout to prove the bound works.
HEALTH_DB_TIMEOUT_SECONDS = 5.0

# Request queue for Gotenberg conversions (P3 §31). Gotenberg's container is
# memory-capped at 1g (docker-compose.yml, P3 §30) on a shared 12GB VM that
# also runs the API, Postgres-adjacent services, and cloudflared — a burst of
# concurrent LibreOffice/Chromium conversions can each spike well past a
# couple hundred MB, so an unbounded number of in-flight requests can OOM the
# container under load. The report's own suggestion (BullMQ + Redis) doesn't
# fit: this is a FastAPI/asyncio codebase with no Node runtime and no Redis
# anywhere else in the stack (rate limiting in middleware.py is already a
# plain in-process structure for the same reason). An ``asyncio.Semaphore``
# is the equivalent in-process primitive — same "acceptable per-worker, not
# cross-worker" caveat as the rate limiter (uvicorn runs 4 workers, so the
# real ceiling is up to ~4x this per instance), which is fine here for the
# same reason: this is defense against a single worker's own burst, not a
# hard global cap.
GOTENBERG_MAX_CONCURRENT = int(os.getenv("GOTENBERG_MAX_CONCURRENT", "2"))
# How long a request waits in queue for a free slot before giving up. Bounded
# so a pile-up of requests behind two slow conversions fails fast with a
# retryable error instead of each queuing request holding its connection open
# indefinitely (and, upstream, holding open the client's XHR).
GOTENBERG_QUEUE_TIMEOUT_SECONDS = 30.0

_gotenberg_semaphore = asyncio.Semaphore(GOTENBERG_MAX_CONCURRENT)


class ConversionQueueTimeout(Exception):
    """Raised when a request waited GOTENBERG_QUEUE_TIMEOUT_SECONDS for a free
    Gotenberg slot and none opened up."""


metrics = {
    "conversions": defaultdict(int),
    "failures": defaultdict(int),
    "total_duration_ms": defaultdict(float),
}


def _error_response(
    msg: str, error_type: str, status: int = 400, headers: dict | None = None
) -> Response:
    return Response(
        content=json.dumps({"error": msg, "error_type": error_type}),
        status_code=status,
        media_type="application/json",
        headers=headers,
    )


def _record_metric(tool_id: str, success: bool, duration_ms: float):
    if success:
        metrics["conversions"][tool_id] += 1
    else:
        metrics["failures"][tool_id] += 1
    metrics["total_duration_ms"][tool_id] += duration_ms


async def _gotenberg_request(endpoint: str, files: dict, tool_id: str) -> bytes:
    url = f"{GOTENBERG_URL}{endpoint}"
    rid = request_id_var.get("-")
    headers = {"X-Request-Id": rid}

    # Queue behind GOTENBERG_MAX_CONCURRENT other in-flight Gotenberg calls
    # (P3 §31) rather than firing straight through — see the module-level
    # comment on _gotenberg_semaphore for why this exists and why it's a
    # semaphore rather than an external queue.
    queue_start = time.time()
    try:
        await asyncio.wait_for(
            _gotenberg_semaphore.acquire(), timeout=GOTENBERG_QUEUE_TIMEOUT_SECONDS
        )
    except TimeoutError:
        queue_ms = round((time.time() - queue_start) * 1000, 1)
        logger.warning(
            "Gotenberg queue timeout: %s waited %sms for a free slot",
            tool_id,
            queue_ms,
            extra={
                "data": {
                    "event": "gotenberg_queue_timeout",
                    "tool_id": tool_id,
                    "gotenberg_endpoint": endpoint,
                    "queue_ms": queue_ms,
                }
            },
        )
        raise ConversionQueueTimeout(
            "The conversion service is busy right now. Please try again in a moment."
        ) from None

    queue_ms = round((time.time() - queue_start) * 1000, 1)
    try:
        start = time.time()
        async with httpx.AsyncClient(timeout=REQUEST_TIMEOUT) as client:
            resp = await client.post(url, files=files, headers=headers)
        gotenberg_ms = round((time.time() - start) * 1000, 1)
    finally:
        _gotenberg_semaphore.release()

    if queue_ms > 50:
        # Only log queueing that actually happened — most requests acquire
        # the semaphore immediately and this would otherwise fire on every
        # single conversion.
        logger.info(
            "Gotenberg request queued %sms before starting: %s",
            queue_ms,
            tool_id,
            extra={
                "data": {
                    "event": "gotenberg_queued",
                    "tool_id": tool_id,
                    "gotenberg_endpoint": endpoint,
                    "queue_ms": queue_ms,
                }
            },
        )

    if resp.status_code != 200:
        error_body = resp.text[:500]
        logger.error(
            "Gotenberg error: %s returned %s",
            endpoint,
            resp.status_code,
            extra={
                "data": {
                    "event": "gotenberg_error",
                    "tool_id": tool_id,
                    "gotenberg_endpoint": endpoint,
                    "gotenberg_status": resp.status_code,
                    "gotenberg_body": error_body,
                    "gotenberg_ms": gotenberg_ms,
                }
            },
        )
        raise RuntimeError(
            f"Gotenberg {endpoint} returned {resp.status_code}: {error_body}"
        )

    logger.info(
        "Gotenberg OK: %s %sms",
        endpoint,
        gotenberg_ms,
        extra={
            "data": {
                "event": "gotenberg_call",
                "tool_id": tool_id,
                "gotenberg_endpoint": endpoint,
                "gotenberg_status": 200,
                "gotenberg_ms": gotenberg_ms,
                "gotenberg_output_bytes": len(resp.content),
            }
        },
    )
    return resp.content


async def _convert_libreoffice(
    content: bytes,
    filename: str,
    extra_form: dict | None = None,
) -> bytes:
    files = {"files": (filename, content)}
    if extra_form:
        for k, v in extra_form.items():
            files[k] = (None, v)
    return await _gotenberg_request(
        "/forms/libreoffice/convert",
        files,
        "libreoffice",
    )


async def _convert_chromium_html(content: bytes, filename: str) -> bytes:
    return await _gotenberg_request(
        "/forms/chromium/convert/html",
        {"files": ("index.html", content)},
        "chromium",
    )


# EPUB -> PDF genuinely does NOT reuse _convert_libreoffice, despite this
# codebase's original assumption that it would: confirmed live against
# production (two independent, spec-compliant EPUB fixtures — one hand-built,
# one generated by the well-known `ebooklib` library — both failed identically
# through Gotenberg's /forms/libreoffice/convert route). LibreOffice has no
# built-in EPUB *import* filter in a stock install (only export), so Gotenberg
# doesn't accept .epub on that route at all. What actually works, and is what
# this uses instead: parse the EPUB's own structure (a ZIP of XHTML chapters
# plus an OPF manifest/spine, both open standards) with the standard library
# only, flatten every chapter into one HTML document in reading order, and
# feed that through the *other* Gotenberg engine already proven in production
# for html-to-pdf — _convert_chromium_html above. Images/CSS are inlined as
# data URIs / <style> blocks rather than sent to Gotenberg as separate
# attached files, since that's guaranteed to resolve regardless of exactly
# how Gotenberg's multi-file placement semantics work for this route.
_EPUB_CONTAINER_NS = {"c": "urn:oasis:names:tc:opendocument:xmlns:container"}
_EPUB_OPF_NS = {"o": "http://www.idpf.org/2007/opf"}
_EPUB_BODY_RE = re.compile(r"<body[^>]*>(.*)</body>", re.IGNORECASE | re.DOTALL)
_EPUB_CSS_LINK_RE = re.compile(
    r'<link[^>]+rel=["\']stylesheet["\'][^>]*href=["\']([^"\']+)["\']', re.IGNORECASE
)
_EPUB_IMG_SRC_RE = re.compile(r'(<img[^>]+src=)(["\'])([^"\']+)\2', re.IGNORECASE)


def _epub_read_text(zf, path: str) -> str:
    raw = zf.read(path)
    return raw.decode("utf-8", errors="replace")


def _epub_resolve(base_dir: str, href: str) -> str:
    return posixpath.normpath(posixpath.join(base_dir, href))


def _epub_data_uri(zf, names: set, path: str) -> str | None:
    if path not in names:
        return None
    mime, _ = mimetypes.guess_type(path)
    data = zf.read(path)
    return f"data:{mime or 'application/octet-stream'};base64,{base64.b64encode(data).decode('ascii')}"


def _safe_xml_fromstring(data: bytes) -> ET.Element:
    """The only server-side parsing of arbitrary uploaded XML in this
    codebase (every other XML tool here runs client-side via the browser's
    DOMParser) — stdlib ElementTree doesn't fully guard against internal
    entity-expansion ("billion laughs") DoS on its own. A real EPUB
    container.xml/OPF manifest never legitimately declares a DOCTYPE, so
    rejecting one outright blocks both that and the external-entity surface
    at once, without pulling in defusedxml for this one narrow case.
    """
    if b"<!DOCTYPE" in data.upper():
        raise ValueError(
            "EPUB control file declares a DOCTYPE, which real container.xml/OPF files never do"
        )
    return ET.fromstring(data)


EPUB_MAX_UNCOMPRESSED_BYTES = 200 * 1024 * 1024  # 200MB, well above any real ebook


def _flatten_epub_to_html_sync(content: bytes) -> bytes:
    with zipfile.ZipFile(io.BytesIO(content)) as zf:
        # Bound the actual decompressed cost, not just the (already-capped)
        # uploaded byte count — the same reasoning as PNG_TO_SVG_MAX_DIMENSION
        # and PDF_TO_PPTX_MAX_RASTER_DIMENSION elsewhere in this file: a small
        # zip can still be crafted to decompress to a huge size ("zip bomb"),
        # and every member gets read into memory below. infolist() exposes
        # each entry's uncompressed size without decompressing it.
        total_uncompressed = sum(zi.file_size for zi in zf.infolist())
        if total_uncompressed > EPUB_MAX_UNCOMPRESSED_BYTES:
            raise ValueError(
                f"EPUB decompresses to {total_uncompressed} bytes, "
                f"over the {EPUB_MAX_UNCOMPRESSED_BYTES} limit"
            )

        names = set(zf.namelist())

        container_root = _safe_xml_fromstring(zf.read("META-INF/container.xml"))
        rootfile = container_root.find(".//c:rootfile", _EPUB_CONTAINER_NS)
        opf_path = rootfile.attrib["full-path"]
        opf_dir = posixpath.dirname(opf_path)

        opf_root = _safe_xml_fromstring(zf.read(opf_path))
        manifest = {
            item.attrib["id"]: item.attrib["href"]
            for item in opf_root.findall(".//o:manifest/o:item", _EPUB_OPF_NS)
        }
        spine_ids = [
            itemref.attrib["idref"]
            for itemref in opf_root.findall(".//o:spine/o:itemref", _EPUB_OPF_NS)
        ]

        style_blocks = []
        seen_css: set[str] = set()
        body_parts = []

        for chapter_num, item_id in enumerate(spine_ids):
            href = manifest.get(item_id)
            if not href:
                continue
            chapter_path = _epub_resolve(opf_dir, href)
            if chapter_path not in names:
                continue
            chapter_dir = posixpath.dirname(chapter_path)
            text = _epub_read_text(zf, chapter_path)

            for css_href in _EPUB_CSS_LINK_RE.findall(text):
                css_path = _epub_resolve(chapter_dir, css_href)
                if css_path in names and css_path not in seen_css:
                    seen_css.add(css_path)
                    style_blocks.append(_epub_read_text(zf, css_path))

            def _inline_img(m, chapter_dir=chapter_dir):
                img_path = _epub_resolve(chapter_dir, m.group(3))
                data_uri = _epub_data_uri(zf, names, img_path)
                return f"{m.group(1)}{m.group(2)}{data_uri or m.group(3)}{m.group(2)}"

            text = _EPUB_IMG_SRC_RE.sub(_inline_img, text)

            body_match = _EPUB_BODY_RE.search(text)
            inner = body_match.group(1) if body_match else text
            # Every chapter but the first starts on a new PDF page.
            page_break = "" if chapter_num == 0 else 'style="page-break-before:always"'
            body_parts.append(f"<section {page_break}>{inner}</section>")

        html = (
            "<!DOCTYPE html><html><head><meta charset=\"utf-8\">"
            f"<style>{''.join(style_blocks)}</style></head><body>"
            f"{''.join(body_parts)}</body></html>"
        )
        return html.encode("utf-8")


async def _convert_epub_to_pdf(content: bytes, filename: str) -> bytes:
    loop = asyncio.get_running_loop()
    html_bytes = await loop.run_in_executor(None, _flatten_epub_to_html_sync, content)
    return await _convert_chromium_html(html_bytes, "index.html")


def _compress_ghostscript_sync(content: bytes, quality: str) -> bytes:
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_in:
        tmp_in.write(content)
        tmp_in_path = tmp_in.name
    tmp_out_path = tmp_in_path + ".out.pdf"

    try:
        start = time.time()
        result = subprocess.run(
            [
                "gs",
                "-sDEVICE=pdfwrite",
                "-dCompatibilityLevel=1.4",
                f"-dPDFSETTINGS=/{quality}",
                "-dNOPAUSE",
                "-dQUIET",
                "-dBATCH",
                f"-sOutputFile={tmp_out_path}",
                tmp_in_path,
            ],
            capture_output=True,
            timeout=REQUEST_TIMEOUT,
        )
        gs_ms = round((time.time() - start) * 1000, 1)

        if result.returncode != 0:
            stderr = result.stderr.decode("utf-8", errors="replace")[:500]
            logger.error(
                "Ghostscript failed: exit %s",
                result.returncode,
                extra={
                    "data": {
                        "event": "ghostscript_error",
                        "gs_returncode": result.returncode,
                        "gs_stderr": stderr,
                        "gs_quality": quality,
                        "gs_ms": gs_ms,
                    }
                },
            )
            raise RuntimeError(f"Ghostscript exit {result.returncode}: {stderr}")

        with open(tmp_out_path, "rb") as f:
            output = f.read()

        logger.info(
            "Ghostscript OK: %s %sms",
            quality,
            gs_ms,
            extra={
                "data": {
                    "event": "ghostscript_call",
                    "gs_quality": quality,
                    "gs_ms": gs_ms,
                    "gs_input_bytes": len(content),
                    "gs_output_bytes": len(output),
                }
            },
        )
        return output
    finally:
        for p in (tmp_in_path, tmp_out_path):
            try:
                os.unlink(p)
            except OSError:
                pass


async def _compress_ghostscript(content: bytes, quality: str) -> bytes:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(
        None, _compress_ghostscript_sync, content, quality
    )


def _convert_pdf2docx_sync(content: bytes) -> bytes:
    from pdf2docx import Converter

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_in:
        tmp_in.write(content)
        tmp_in_path = tmp_in.name
    tmp_out_path = tmp_in_path + ".docx"

    try:
        start = time.time()
        cv = Converter(tmp_in_path)
        cv.convert(tmp_out_path)
        cv.close()
        p2d_ms = round((time.time() - start) * 1000, 1)

        with open(tmp_out_path, "rb") as f:
            output = f.read()

        logger.info(
            "pdf2docx OK: %sms",
            p2d_ms,
            extra={
                "data": {
                    "event": "pdf2docx_call",
                    "p2d_ms": p2d_ms,
                    "p2d_input_bytes": len(content),
                    "p2d_output_bytes": len(output),
                }
            },
        )
        return output
    except Exception:
        logger.error(
            "pdf2docx failed",
            exc_info=True,
            extra={
                "data": {
                    "event": "pdf2docx_error",
                    "p2d_input_bytes": len(content),
                }
            },
        )
        raise
    finally:
        for p in (tmp_in_path, tmp_out_path):
            try:
                os.unlink(p)
            except OSError:
                pass


async def _convert_pdf2docx(content: bytes, filename: str) -> bytes:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _convert_pdf2docx_sync, content)


# pdfplumber's default table-detection ("lines" strategy) only finds tables
# with real ruled/bordered grid lines. Most real-world "tables" people
# actually want extracted — bank statements, invoices, whitespace-aligned
# reports — have no ruling at all, so the lines strategy finds nothing on
# them and every row would otherwise collapse into the plain-text fallback
# (one unsplit line per row, exactly the "not really a spreadsheet" outcome
# this tool exists to avoid). Tried second, only when the lines strategy
# finds nothing, since the text strategy has its own false-positive risk
# (it can mis-split large standalone headings into spurious "columns").
_PDF_TO_XLSX_TEXT_TABLE_SETTINGS = {
    "vertical_strategy": "text",
    "horizontal_strategy": "text",
}

# openpyxl auto-detects any string value starting with "=" as a live Excel
# formula (confirmed directly: Cell._bind_value sets data_type="f", which
# serializes as a real <f> element) — so writing untrusted PDF text straight
# into a cell lets a crafted PDF plant a formula (e.g. =WEBSERVICE(...) or
# =HYPERLINK(...)) that runs the moment the converted file is opened in
# Excel, no macro warning involved (CWE-1236).
#
# Deliberately just "=", not the usual CSV-injection "=+-@" prefix set:
# verified directly that openpyxl does NOT auto-formula-type "+"/"-"/"@"
# (only "=" flips data_type to "f") — and "-" is the single most common
# leading character in exactly the data this tool extracts (negative
# amounts on a bank statement or invoice), so blanket-sanitizing it would
# silently corrupt real numbers into text on every debit line. _XLSX_NUMERIC_RE
# below also means genuine numbers never reach this sanitizer as strings at
# all — they're converted to a real int/float first, sidestepping the
# question entirely for the common case.
_XLSX_FORMULA_TRIGGER_CHARS = ("=",)

_XLSX_NUMERIC_RE = re.compile(r"^[+-]?\d+(\.\d+)?$")


def _sanitize_xlsx_cell(value):
    if isinstance(value, str) and value.startswith(_XLSX_FORMULA_TRIGGER_CHARS):
        return "'" + value
    return value


def _coerce_xlsx_numeric(value: str):
    """A number extracted as text ("-4.50", "1,234") stays text in the
    output cell unless converted — and a text cell doesn't participate in
    Excel's own SUM/sort/filter the way this tool's own content promises
    ("Formulas, sums, sorting, filtering"). Narrow on purpose: plain digits,
    an optional leading sign, an optional single decimal point, optional
    comma thousands separators — not currency symbols or accounting-style
    parenthesized negatives, which would mean guessing at ambiguous formats
    rather than a real fix.
    """
    stripped = value.strip().replace(",", "")
    if not _XLSX_NUMERIC_RE.match(stripped):
        return value
    try:
        return int(stripped) if "." not in stripped else float(stripped)
    except ValueError:
        return value


def _prepare_xlsx_cell_value(value):
    if isinstance(value, str):
        numeric = _coerce_xlsx_numeric(value)
        if not isinstance(numeric, str):
            return numeric
        return _sanitize_xlsx_cell(value)
    return value


def _convert_pdf_to_xlsx_sync(content: bytes) -> bytes:
    """PDF -> XLSX. No Gotenberg/LibreOffice route does this (LibreOffice's
    PDF import produces a Draw/Writer document, not spreadsheet cells), so
    this is a pure-Python step in the same shape as pdf2docx above: read
    in-memory (pdfplumber/openpyxl both support file-like objects directly,
    so no temp files are needed here), one sheet per PDF page. A page with a
    detected table gets real cells; a page without one still gets its plain
    extracted text (one line per row) rather than being silently dropped —
    honest about pages that don't contain a real table, not a failure.
    """
    import pdfplumber
    from openpyxl import Workbook

    start = time.time()
    try:
        wb = Workbook()
        wb.remove(wb.active)

        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                ws = wb.create_sheet(title=f"Page {page_num}")
                tables = page.extract_tables()
                if not tables:
                    tables = page.extract_tables(
                        table_settings=_PDF_TO_XLSX_TEXT_TABLE_SETTINGS
                    )
                if tables:
                    row_cursor = 1
                    for table in tables:
                        for row in table:
                            # The text strategy emits fully-blank spacer rows
                            # between detected lines (gaps in the page it
                            # couldn't assign to a real row) — skip those
                            # rather than writing empty rows into the sheet.
                            if not any(cell for cell in row):
                                continue
                            for col_idx, cell in enumerate(row, start=1):
                                ws.cell(
                                    row=row_cursor,
                                    column=col_idx,
                                    value=_prepare_xlsx_cell_value(cell),
                                )
                            row_cursor += 1
                        row_cursor += 1  # blank row between tables on the same page
                else:
                    text = page.extract_text() or ""
                    for row_idx, line in enumerate(text.splitlines(), start=1):
                        ws.cell(
                            row=row_idx, column=1, value=_prepare_xlsx_cell_value(line)
                        )

        if not wb.sheetnames:
            wb.create_sheet(title="Sheet1")

        buf = io.BytesIO()
        wb.save(buf)
        output = buf.getvalue()
        p2x_ms = round((time.time() - start) * 1000, 1)
        logger.info(
            "pdf-to-xlsx OK: %sms",
            p2x_ms,
            extra={
                "data": {
                    "event": "pdf_to_xlsx_call",
                    "p2x_ms": p2x_ms,
                    "p2x_input_bytes": len(content),
                    "p2x_output_bytes": len(output),
                }
            },
        )
        return output
    except Exception:
        logger.error(
            "pdf-to-xlsx failed",
            exc_info=True,
            extra={
                "data": {"event": "pdf_to_xlsx_error", "p2x_input_bytes": len(content)}
            },
        )
        raise


async def _convert_pdf_to_xlsx(content: bytes, filename: str) -> bytes:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _convert_pdf_to_xlsx_sync, content)


# PDF's MediaBox is spec-legal up to 14,400x14,400pt (200x200in), and a
# near-blank page declaring that size costs almost nothing in file bytes —
# confirmed directly: a single blank 3000x3000pt page is a 520-byte PDF, but
# rendering it at a flat 150 DPI with no cap allocates a 112MB pixmap in
# ~20ms; scaling to the spec maximum extrapolates to gigabytes for one page.
# This runs inside the same shared thread-pool executor every other
# conversion route uses (loop.run_in_executor(None, ...)), so an unbounded
# render can OOM the whole API process, not just this one request. Mirrors
# PNG_TO_SVG_MAX_DIMENSION's role for the tracer below: bound the actual
# rendered-pixel cost per page, independent of the page's declared size —
# not just the already-capped upload byte count, which this cheaply defeats.
PDF_TO_PPTX_TARGET_DPI = 150
PDF_TO_PPTX_MAX_RASTER_DIMENSION = 3000  # px, longest side of any one page's render

# PowerPoint's own practical maximum custom slide size (each axis). Without
# this, an unusually large but entirely legitimate page — a wide
# engineering-drawing PDF, say — could produce a .pptx PowerPoint itself
# refuses to open, independent of anything hostile.
PPTX_MAX_SLIDE_EMU = 56 * 914400  # 56in, per axis


def _convert_pdf_to_pptx_sync(content: bytes) -> bytes:
    """PDF -> PPTX. Same "no Gotenberg route for this direction" reasoning as
    XLSX above. Renders each page to an image (PyMuPDF) and places it on its
    own slide (python-pptx) — preserves the page's exact visual appearance,
    at the cost of the slide's text not being separately editable, which is
    disclosed in the tool's own FAQ rather than left as a surprise. The deck
    is sized to the first page's aspect ratio; any page with a different
    aspect ratio is letterboxed (not stretched) to avoid distorting it.
    """
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Emu

    start = time.time()
    try:
        doc = fitz.open(stream=content, filetype="pdf")
        try:
            if doc.page_count == 0:
                raise RuntimeError("PDF has no pages")

            prs = Presentation()
            first_rect = doc[0].rect
            prs.slide_width = Emu(
                min(int(first_rect.width / 72 * 914400), PPTX_MAX_SLIDE_EMU)
            )
            prs.slide_height = Emu(
                min(int(first_rect.height / 72 * 914400), PPTX_MAX_SLIDE_EMU)
            )
            blank_layout = prs.slide_layouts[6]

            for page in doc:
                # Target DPI first, but never let the rendered raster exceed
                # the pixel-dimension cap regardless of this page's own
                # declared size — independent of the slide EMU size above,
                # since the image is scaled to fit the slide either way.
                zoom = PDF_TO_PPTX_TARGET_DPI / 72
                longest_pt = max(page.rect.width, page.rect.height)
                if longest_pt * zoom > PDF_TO_PPTX_MAX_RASTER_DIMENSION:
                    zoom = PDF_TO_PPTX_MAX_RASTER_DIMENSION / longest_pt
                pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                img_bytes = pix.tobytes("png")

                img_aspect = pix.width / pix.height
                slide_aspect = prs.slide_width / prs.slide_height
                if img_aspect > slide_aspect:
                    draw_w = prs.slide_width
                    draw_h = int(draw_w / img_aspect)
                else:
                    draw_h = prs.slide_height
                    draw_w = int(draw_h * img_aspect)
                left = (prs.slide_width - draw_w) // 2
                top = (prs.slide_height - draw_h) // 2

                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    io.BytesIO(img_bytes), left, top, width=draw_w, height=draw_h
                )
        finally:
            doc.close()

        buf = io.BytesIO()
        prs.save(buf)
        output = buf.getvalue()
        p2p_ms = round((time.time() - start) * 1000, 1)
        logger.info(
            "pdf-to-pptx OK: %sms",
            p2p_ms,
            extra={
                "data": {
                    "event": "pdf_to_pptx_call",
                    "p2p_ms": p2p_ms,
                    "p2p_input_bytes": len(content),
                    "p2p_output_bytes": len(output),
                }
            },
        )
        return output
    except Exception:
        logger.error(
            "pdf-to-pptx failed",
            exc_info=True,
            extra={
                "data": {"event": "pdf_to_pptx_error", "p2p_input_bytes": len(content)}
            },
        )
        raise


async def _convert_pdf_to_pptx(content: bytes, filename: str) -> bytes:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _convert_pdf_to_pptx_sync, content)


# PNG to SVG (Vectorize) — a genuinely new server-side capability, not a
# reuse of the Gotenberg/LibreOffice document pipeline the other 3 Cloud
# tools in this PR build on. No potrace/vtracer dependency: this hand-rolls
# color-quantize + contour-trace + polygon-simplify on top of Pillow/OpenCV,
# the same "own the well-specified algorithm, lean on a vetted library for
# the primitives" approach this codebase already uses for MD5, RC4, Code 128,
# QR, and the ICO container — chosen over potrace (monochrome-only output)
# and vtracer (a Rust-toolchain risk for this image, with no guaranteed
# prebuilt wheel for this target).
PNG_TO_SVG_MAX_DIMENSION = 1500  # px, longest side, after decode
PNG_TO_SVG_PALETTE_SIZE = 16
PNG_TO_SVG_MIN_CONTOUR_AREA = 6  # px^2 — drops single/few-pixel noise fragments


def _contour_to_svg_path_d(contour) -> str:
    import cv2

    # 1.0px simplification tolerance: smooths jagged pixel-grid edges from
    # the mask boundary without losing the shape's real corners.
    approx = cv2.approxPolyDP(contour, 1.0, True)
    points = approx.reshape(-1, 2)
    if len(points) < 3:
        return ""
    d = f"M {points[0][0]} {points[0][1]} "
    d += " ".join(f"L {x} {y}" for x, y in points[1:])
    d += " Z"
    return d


def _trace_png_to_svg_sync(content: bytes) -> bytes:
    import cv2
    import numpy as np
    from PIL import Image

    start = time.time()
    try:
        img = Image.open(io.BytesIO(content)).convert("RGBA")

        # Bound tracing cost independent of the (already size-capped) upload:
        # a small, heavily-compressed PNG can still decode to a huge pixel grid.
        w, h = img.size
        if max(w, h) > PNG_TO_SVG_MAX_DIMENSION:
            scale = PNG_TO_SVG_MAX_DIMENSION / max(w, h)
            img = img.resize(
                (max(1, round(w * scale)), max(1, round(h * scale))), Image.LANCZOS
            )
            w, h = img.size

        # This tracer's SVG output has no alpha channel, so flatten
        # transparency onto a white background rather than leave it
        # undefined per-pixel.
        background = Image.new("RGBA", img.size, (255, 255, 255, 255))
        flattened = Image.alpha_composite(background, img).convert("RGB")

        quantized = flattened.quantize(
            colors=PNG_TO_SVG_PALETTE_SIZE, method=Image.MEDIANCUT
        )
        palette = quantized.getpalette()
        indexed = np.array(quantized)

        paths: list[str] = []
        for color_index in range(PNG_TO_SVG_PALETTE_SIZE):
            mask = (indexed == color_index).astype(np.uint8)
            if not mask.any():
                continue
            r, g, b = palette[color_index * 3 : color_index * 3 + 3]
            hex_color = f"#{r:02x}{g:02x}{b:02x}"

            contours, hierarchy = cv2.findContours(
                mask * 255, cv2.RETR_CCOMP, cv2.CHAIN_APPROX_SIMPLE
            )
            if hierarchy is None:
                continue
            hierarchy = hierarchy[0]

            # RETR_CCOMP gives each outer contour and its direct hole
            # children (e.g. the inside of an "O") as separate entries —
            # group them so one <path fill-rule="evenodd"> per outer shape
            # renders holes correctly instead of filling them in solid.
            children_by_parent: dict[int, list[int]] = {}
            top_level: list[int] = []
            for i, (_next, _prev, _child, parent) in enumerate(hierarchy):
                if parent == -1:
                    top_level.append(i)
                else:
                    children_by_parent.setdefault(parent, []).append(i)

            for outer_idx in top_level:
                if cv2.contourArea(contours[outer_idx]) < PNG_TO_SVG_MIN_CONTOUR_AREA:
                    continue
                subpaths = [_contour_to_svg_path_d(contours[outer_idx])]
                for hole_idx in children_by_parent.get(outer_idx, []):
                    subpaths.append(_contour_to_svg_path_d(contours[hole_idx]))
                d = " ".join(p for p in subpaths if p)
                if d:
                    paths.append(
                        f'<path d="{d}" fill="{hex_color}" fill-rule="evenodd"/>'
                    )

        svg = (
            f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {w} {h}" '
            f'width="{w}" height="{h}">\n' + "\n".join(paths) + "\n</svg>\n"
        )
        output = svg.encode("utf-8")
        trace_ms = round((time.time() - start) * 1000, 1)
        logger.info(
            "png-to-svg OK: %sms",
            trace_ms,
            extra={
                "data": {
                    "event": "png_to_svg_call",
                    "trace_ms": trace_ms,
                    "trace_input_bytes": len(content),
                    "trace_output_bytes": len(output),
                    "trace_path_count": len(paths),
                }
            },
        )
        return output
    except Exception:
        logger.error(
            "png-to-svg failed",
            exc_info=True,
            extra={
                "data": {"event": "png_to_svg_error", "trace_input_bytes": len(content)}
            },
        )
        raise


async def _trace_png_to_svg(content: bytes, filename: str) -> bytes:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _trace_png_to_svg_sync, content)


def _safe_filename(filename: str, ext: str = ".pdf") -> str:
    base = filename.rsplit(".", 1)[0] if "." in filename else filename
    safe = "".join(c for c in base if c.isalnum() or c in " ._-")
    return (safe.strip() or "converted") + ext


def _max_size_for(user: User | None) -> int:
    """Signed-in users get a flat ×2 (spec §6.3); anonymous gets the base cap.
    Server-side deliberately ignores ``user.max_file_size`` to bound the
    expensive Gotenberg/Ghostscript path and match the client rule (§6.1)."""
    return MAX_FILE_SIZE * (2 if user else 1)


_READ_CHUNK_SIZE = 1024 * 1024  # 1MB


async def _read_capped(file: UploadFile, max_size: int) -> bytes:
    """Read ``file`` in chunks, aborting as soon as it exceeds ``max_size``.

    A plain ``await file.read()`` buffers the entire body into memory before
    ``validate_upload``'s size check ever runs — a client can make the API
    fully materialize an arbitrarily large upload as ``bytes`` regardless of
    what it declared. Chunked reads with an early abort mean a request that's
    already over the limit gets rejected without paying for the rest of it.
    """
    chunks: list[bytes] = []
    total = 0
    while True:
        chunk = await file.read(_READ_CHUNK_SIZE)
        if not chunk:
            break
        total += len(chunk)
        if total > max_size:
            raise ValidationError(
                f"File exceeds the {max_size // (1024 * 1024)}MB limit "
                "for this conversion type.",
                "too_large",
                bytes_read=total,
            )
        chunks.append(chunk)
    return b"".join(chunks)


async def _handle_conversion(
    file: UploadFile,
    tool_id: str,
    convert_fn,
    output_ext: str = ".pdf",
    output_mime: str = "application/pdf",
    max_size: int = MAX_FILE_SIZE,
) -> Response:
    start = time.time()
    input_size = 0
    try:
        content = await _read_capped(file, max_size)
        input_size = len(content)
        validate_upload(content, file.filename or "unknown", tool_id, max_size)
        result = await convert_fn(content, file.filename or "file")
        duration_ms = round((time.time() - start) * 1000, 1)
        _record_metric(tool_id, True, duration_ms)

        logger.info(
            "Conversion OK: %s",
            tool_id,
            extra={
                "data": {
                    "event": "conversion_success",
                    "tool_id": tool_id,
                    "input_bytes": input_size,
                    "output_bytes": len(result),
                    "duration_ms": duration_ms,
                }
            },
        )

        output_filename = _safe_filename(file.filename or "file", output_ext)
        return Response(
            content=result,
            media_type=output_mime,
            headers={
                "Content-Disposition": f'attachment; filename="{output_filename}"',
                "X-Conversion-Time": f"{duration_ms}ms",
            },
        )
    except ValidationError as e:
        if e.bytes_read is not None:
            input_size = e.bytes_read
        duration_ms = round((time.time() - start) * 1000, 1)
        _record_metric(tool_id, False, duration_ms)
        logger.warning(
            "Validation failed: %s — %s",
            tool_id,
            e.error_type,
            extra={
                "data": {
                    "event": "validation_error",
                    "tool_id": tool_id,
                    "error_type": e.error_type,
                    "input_bytes": input_size,
                    "duration_ms": duration_ms,
                }
            },
        )
        return _error_response(e.message, e.error_type, 400)
    except ConversionQueueTimeout as e:
        duration_ms = round((time.time() - start) * 1000, 1)
        _record_metric(tool_id, False, duration_ms)
        logger.warning(
            "Conversion queue timeout: %s",
            tool_id,
            extra={
                "data": {
                    "event": "queue_timeout",
                    "tool_id": tool_id,
                    "input_bytes": input_size,
                    "duration_ms": duration_ms,
                }
            },
        )
        # Retryable, so callers get a Retry-After the same way the rate
        # limiter's 429 does (middleware.py) rather than a bare 503.
        return _error_response(
            str(e), "queue_timeout", 503, headers={"Retry-After": "10"}
        )
    except subprocess.TimeoutExpired:
        duration_ms = round((time.time() - start) * 1000, 1)
        _record_metric(tool_id, False, duration_ms)
        logger.error(
            "Conversion timeout: %s",
            tool_id,
            extra={
                "data": {
                    "event": "conversion_timeout",
                    "tool_id": tool_id,
                    "input_bytes": input_size,
                    "duration_ms": duration_ms,
                }
            },
        )
        return _error_response(
            "Conversion timed out. Try a simpler or smaller file.",
            "timeout",
            504,
        )
    except Exception:
        duration_ms = round((time.time() - start) * 1000, 1)
        _record_metric(tool_id, False, duration_ms)
        logger.error(
            "Conversion failed: %s",
            tool_id,
            exc_info=True,
            extra={
                "data": {
                    "event": "conversion_error",
                    "tool_id": tool_id,
                    "input_bytes": input_size,
                    "duration_ms": duration_ms,
                }
            },
        )
        return _error_response(
            "Conversion failed. The file may be corrupted or password-protected.",
            "conversion_error",
            500,
        )


@router.post("/convert/docx-to-pdf")
async def docx_to_pdf(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    return await _handle_conversion(
        file, "docx-to-pdf", _convert_libreoffice, max_size=_max_size_for(user)
    )


@router.post("/convert/xlsx-to-pdf")
async def xlsx_to_pdf(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    async def _convert(content: bytes, filename: str) -> bytes:
        return await _convert_libreoffice(
            content,
            filename,
            extra_form={"landscape": "true", "singlePageSheets": "true"},
        )

    return await _handle_conversion(
        file, "xlsx-to-pdf", _convert, max_size=_max_size_for(user)
    )


@router.post("/convert/pptx-to-pdf")
async def pptx_to_pdf(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    return await _handle_conversion(
        file, "pptx-to-pdf", _convert_libreoffice, max_size=_max_size_for(user)
    )


@router.post("/convert/html-to-pdf")
async def html_to_pdf(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    return await _handle_conversion(
        file, "html-to-pdf", _convert_chromium_html, max_size=_max_size_for(user)
    )


@router.post("/convert/pdf-compress")
async def pdf_compress(
    file: UploadFile = File(...),
    quality: str = Form("ebook"),
    user: User | None = Depends(current_user_for_convert),
):
    quality = validate_compress_quality(quality)

    async def _compress(content: bytes, filename: str) -> bytes:
        return await _compress_ghostscript(content, quality)

    return await _handle_conversion(
        file, "pdf-compress", _compress, max_size=_max_size_for(user)
    )


@router.post("/convert/pdf-to-docx")
async def pdf_to_docx(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    return await _handle_conversion(
        file,
        "pdf-to-docx",
        _convert_pdf2docx,
        output_ext=".docx",
        output_mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        max_size=_max_size_for(user),
    )


@router.post("/convert/pdf-to-xlsx")
async def pdf_to_xlsx(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    return await _handle_conversion(
        file,
        "pdf-to-xlsx",
        _convert_pdf_to_xlsx,
        output_ext=".xlsx",
        output_mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        max_size=_max_size_for(user),
    )


@router.post("/convert/pdf-to-pptx")
async def pdf_to_pptx(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    return await _handle_conversion(
        file,
        "pdf-to-pptx",
        _convert_pdf_to_pptx,
        output_ext=".pptx",
        output_mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        max_size=_max_size_for(user),
    )


@router.post("/convert/epub-to-pdf")
async def epub_to_pdf(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    return await _handle_conversion(
        file, "epub-to-pdf", _convert_epub_to_pdf, max_size=_max_size_for(user)
    )


# Smaller than the 25MB every other Cloud tool allows — tracing is more
# CPU-intensive than a document conversion, and PNG_TO_SVG_MAX_DIMENSION
# already bounds per-request cost independent of this cap (see the tracer
# above); this just keeps very large uploads from queuing up in the first
# place. Matches the tool's own YAML max_file_size (client-side check).
#
# Deliberately flat for every user tier, unlike every other Cloud route:
# min(_max_size_for(user), PNG_TO_SVG_MAX_UPLOAD) below means a signed-in
# user's usual ×2 allowance (_max_size_for) is capped right back down to
# this constant. That's intentional, not an accidental side effect of
# min() — the cap here is driven by tracing's CPU/memory cost per byte,
# which doesn't get cheaper because the uploader has an account, unlike the
# other Cloud tools' cap (mostly abuse/queue-fairness, where trusting a
# signed-in user with more headroom makes sense).
PNG_TO_SVG_MAX_UPLOAD = 10 * 1024 * 1024


@router.post("/convert/png-to-svg")
async def png_to_svg(
    file: UploadFile = File(...),
    user: User | None = Depends(current_user_for_convert),
):
    return await _handle_conversion(
        file,
        "png-to-svg",
        _trace_png_to_svg,
        output_ext=".svg",
        output_mime="image/svg+xml",
        max_size=min(_max_size_for(user), PNG_TO_SVG_MAX_UPLOAD),
    )


async def _check_gotenberg() -> bool:
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            resp = await client.get(f"{GOTENBERG_URL}/health")
        return resp.status_code == 200
    except Exception:
        return False


async def _check_db() -> bool:
    # P2 §20 — an uptime monitor hitting /health should see the DB as part of
    # "is the API actually usable", not just the Gotenberg dependency:
    # ratings/history/auth/admin all go through it, same connectivity check
    # main.py's startup lifespan already does, just per-request instead of
    # once at boot.
    #
    # Post-merge audit fix: bounded with the same 5s budget as the Gotenberg
    # check. `async_engine` (data/db.py) sets no `connect_timeout` — only the
    # sync engine does — so an unbounded `.connect()` here could hang past a
    # monitor's polling interval on a network partition that drops packets
    # instead of refusing the connection outright. Unlike the lifespan's
    # ONE-TIME startup check this pattern mirrors, this one is now
    # re-triggered by every external poll, so a hang here is no longer a
    # one-off — it accumulates a stuck task (and a held pool connection) per
    # poll for as long as the partition lasts.
    try:
        async with asyncio.timeout(HEALTH_DB_TIMEOUT_SECONDS):
            async with async_engine.connect() as conn:
                await conn.execute(text("SELECT 1"))
        return True
    except Exception:
        return False


@router.api_route("/health", methods=["GET", "HEAD"])
async def health():
    # Run concurrently, not sequentially — both checks are independently
    # bounded at 5s, but running them one after another let a partial outage
    # (both dependencies slow-but-not-instantly-refused) push total latency to
    # ~10s, right when a monitor most needs a fast "degraded" signal instead
    # of a bare timeout. gather() bounds total latency at max(5s, 5s).
    gotenberg_ok, db_ok = await asyncio.gather(_check_gotenberg(), _check_db())

    status = "healthy" if (gotenberg_ok and db_ok) else "degraded"
    if not gotenberg_ok:
        logger.warning(
            "Health check: Gotenberg unreachable",
            extra={
                "data": {
                    "event": "health_degraded",
                    "gotenberg": "down",
                }
            },
        )
    if not db_ok:
        logger.warning(
            "Health check: database unreachable",
            extra={
                "data": {
                    "event": "health_degraded",
                    "database": "down",
                }
            },
        )
    return {
        "status": status,
        "gotenberg": "up" if gotenberg_ok else "down",
        "database": "up" if db_ok else "down",
    }


@router.get("/metrics")
async def get_metrics(_admin=Depends(require_admin)):
    return {
        "conversions": dict(metrics["conversions"]),
        "failures": dict(metrics["failures"]),
        "avg_duration_ms": {
            k: round(v / max(metrics["conversions"][k] + metrics["failures"][k], 1), 1)
            for k, v in metrics["total_duration_ms"].items()
        },
    }
