"""Integration — pre-existing converter routes.

Validation/error paths need no Gotenberg (validation runs before the proxy call).
The success path mocks the Gotenberg call so no container is required.
"""

import asyncio
import io

import converter


async def test_convert_rejects_wrong_extension(client):
    r = await client.post(
        "/api/v1/convert/docx-to-pdf",
        files={"file": ("notes.txt", b"hello world", "text/plain")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "wrong_format"


async def test_convert_rejects_empty_file(client):
    r = await client.post(
        "/api/v1/convert/docx-to-pdf",
        files={"file": ("empty.docx", b"", "application/octet-stream")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "empty_file"


async def test_convert_success_with_mocked_gotenberg(client, monkeypatch):
    async def fake_libreoffice(content, filename, extra_form=None):
        return b"%PDF-1.4 fake pdf bytes"

    monkeypatch.setattr(converter, "_convert_libreoffice", fake_libreoffice)
    valid_docx = b"PK\x03\x04" + b"\x00" * 200  # passes magic-byte check
    r = await client.post(
        "/api/v1/convert/docx-to-pdf",
        files={
            "file": (
                "report.docx",
                valid_docx,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )
    assert r.status_code == 200
    assert r.headers["content-type"] == "application/pdf"
    assert r.headers["content-disposition"] == 'attachment; filename="report.pdf"'
    assert r.content == b"%PDF-1.4 fake pdf bytes"


async def test_pdf_to_xlsx_rejects_wrong_extension(client):
    r = await client.post(
        "/api/v1/convert/pdf-to-xlsx",
        files={"file": ("notes.txt", b"hello world", "text/plain")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "wrong_format"


async def test_pdf_to_xlsx_rejects_empty_file(client):
    r = await client.post(
        "/api/v1/convert/pdf-to-xlsx",
        files={"file": ("empty.pdf", b"", "application/pdf")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "empty_file"


async def test_pdf_to_xlsx_success_with_mocked_conversion(client, monkeypatch):
    async def fake_convert(content, filename):
        return b"fake xlsx bytes"

    monkeypatch.setattr(converter, "_convert_pdf_to_xlsx", fake_convert)
    valid_pdf = b"%PDF-1.4" + b"\x00" * 200
    r = await client.post(
        "/api/v1/convert/pdf-to-xlsx",
        files={"file": ("report.pdf", valid_pdf, "application/pdf")},
    )
    assert r.status_code == 200
    assert (
        r.headers["content-type"]
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert r.headers["content-disposition"] == 'attachment; filename="report.xlsx"'
    assert r.content == b"fake xlsx bytes"


async def test_pdf_to_pptx_rejects_wrong_extension(client):
    r = await client.post(
        "/api/v1/convert/pdf-to-pptx",
        files={"file": ("notes.txt", b"hello world", "text/plain")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "wrong_format"


async def test_pdf_to_pptx_rejects_empty_file(client):
    r = await client.post(
        "/api/v1/convert/pdf-to-pptx",
        files={"file": ("empty.pdf", b"", "application/pdf")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "empty_file"


async def test_pdf_to_pptx_success_with_mocked_conversion(client, monkeypatch):
    async def fake_convert(content, filename):
        return b"fake pptx bytes"

    monkeypatch.setattr(converter, "_convert_pdf_to_pptx", fake_convert)
    valid_pdf = b"%PDF-1.4" + b"\x00" * 200
    r = await client.post(
        "/api/v1/convert/pdf-to-pptx",
        files={"file": ("deck.pdf", valid_pdf, "application/pdf")},
    )
    assert r.status_code == 200
    assert (
        r.headers["content-type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert r.headers["content-disposition"] == 'attachment; filename="deck.pptx"'
    assert r.content == b"fake pptx bytes"


async def test_epub_to_pdf_rejects_wrong_extension(client):
    r = await client.post(
        "/api/v1/convert/epub-to-pdf",
        files={"file": ("notes.txt", b"hello world", "text/plain")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "wrong_format"


async def test_epub_to_pdf_rejects_empty_file(client):
    r = await client.post(
        "/api/v1/convert/epub-to-pdf",
        files={"file": ("empty.epub", b"", "application/epub+zip")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "empty_file"


async def test_epub_to_pdf_success_with_mocked_gotenberg(client, monkeypatch):
    async def fake_libreoffice(content, filename, extra_form=None):
        return b"%PDF-1.4 fake pdf bytes"

    monkeypatch.setattr(converter, "_convert_libreoffice", fake_libreoffice)
    valid_epub = b"PK\x03\x04" + b"\x00" * 200  # passes magic-byte check
    r = await client.post(
        "/api/v1/convert/epub-to-pdf",
        files={"file": ("book.epub", valid_epub, "application/epub+zip")},
    )
    assert r.status_code == 200
    assert r.headers["content-type"] == "application/pdf"
    assert r.headers["content-disposition"] == 'attachment; filename="book.pdf"'
    assert r.content == b"%PDF-1.4 fake pdf bytes"


async def test_png_to_svg_rejects_wrong_extension(client):
    r = await client.post(
        "/api/v1/convert/png-to-svg",
        files={"file": ("notes.txt", b"hello world", "text/plain")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "wrong_format"


async def test_png_to_svg_rejects_empty_file(client):
    r = await client.post(
        "/api/v1/convert/png-to-svg",
        files={"file": ("empty.png", b"", "image/png")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "empty_file"


async def test_png_to_svg_rejects_oversized_file(client):
    # png-to-svg caps uploads at 10MB, below the 25MB other Cloud tools
    # allow (tracing is more CPU-intensive) — proves that lower cap is
    # actually enforced, not just the shared default.
    oversized = b"\x89PNG\r\n\x1a\n" + b"\x00" * (10 * 1024 * 1024 + 1)
    r = await client.post(
        "/api/v1/convert/png-to-svg",
        files={"file": ("huge.png", oversized, "image/png")},
    )
    assert r.status_code == 400
    assert r.json()["error_type"] == "too_large"


async def test_png_to_svg_success_with_mocked_conversion(client, monkeypatch):
    async def fake_trace(content, filename):
        return b"<svg>fake</svg>"

    monkeypatch.setattr(converter, "_trace_png_to_svg", fake_trace)
    valid_png = b"\x89PNG\r\n\x1a\n" + b"\x00" * 200
    r = await client.post(
        "/api/v1/convert/png-to-svg",
        files={"file": ("logo.png", valid_png, "image/png")},
    )
    assert r.status_code == 200
    assert r.headers["content-type"] == "image/svg+xml"
    assert r.headers["content-disposition"] == 'attachment; filename="logo.svg"'
    assert r.content == b"<svg>fake</svg>"


# --------------------------------------------------------------------------- #
# Direct, unmocked tests of the 3 new conversion functions' actual logic.
#
# The route-level tests above all monkeypatch the conversion function itself,
# so they only prove the HTTP plumbing (validation, headers, error mapping) —
# none of them exercise the real table extraction, slide layout, or tracing
# algorithm. These call the _sync functions directly against real generated
# PDFs/PNGs, mirroring this codebase's existing precedent for a newly-added,
# non-trivial algorithm getting its own direct coverage beyond the route
# wrapper (test_pdf_lib_worker_crypto.test.js's RC4 round trip, the ICO
# round-trip test, etc.).
# --------------------------------------------------------------------------- #


def _make_pdf_with_ruled_table():
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    rows = [["Item", "Qty", "Price"], ["Widget", "10", "$5.00"]]
    x0, y0, col_w, row_h = 72, 100, 120, 24
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            x, y = x0 + c * col_w, y0 + r * row_h
            page.draw_rect(
                fitz.Rect(x, y, x + col_w, y + row_h), color=(0, 0, 0), width=0.5
            )
            page.insert_text((x + 4, y + 16), cell, fontsize=11)
    content = doc.tobytes()
    doc.close()
    return content


def test_pdf_to_xlsx_extracts_ruled_table_into_real_cells():
    from openpyxl import load_workbook

    xlsx_bytes = converter._convert_pdf_to_xlsx_sync(_make_pdf_with_ruled_table())
    wb = load_workbook(io.BytesIO(xlsx_bytes))
    ws = wb["Page 1"]
    rows = [tuple(r) for r in ws.iter_rows(min_row=1, max_row=2, values_only=True)]
    # "10" comes back as a real int, not the string "10" — a number-looking
    # cell must actually be numeric or Excel's own SUM/sort/filter can't use
    # it. "$5.00" stays text: the currency symbol makes it deliberately not
    # auto-parsed (see _coerce_xlsx_numeric's docstring).
    assert rows == [("Item", "Qty", "Price"), ("Widget", 10, "$5.00")]
    assert ws["B2"].data_type == "n"  # real numeric cell, not text


def test_pdf_to_xlsx_falls_back_to_text_when_no_table():
    import fitz
    from openpyxl import load_workbook

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 72), "Just a plain paragraph, no table here.", fontsize=14)
    content = doc.tobytes()
    doc.close()

    xlsx_bytes = converter._convert_pdf_to_xlsx_sync(content)
    wb = load_workbook(io.BytesIO(xlsx_bytes))
    ws = wb["Page 1"]
    assert ws["A1"].value == "Just a plain paragraph, no table here."


def test_pdf_to_xlsx_detects_borderless_whitespace_aligned_table():
    # The realistic case (bank statements, invoices) — no ruled lines, so the
    # default "lines" table strategy alone finds nothing and everything would
    # otherwise collapse into one unsplit text blob per row. 3 rows, not 2:
    # pdfplumber's text strategy needs >= 3 words aligned in a column
    # (min_words_vertical, its own default) before it infers a vertical
    # divider there at all — a 2-row page finds nothing under either
    # strategy, which is a real, inherent floor of this approach, not
    # something this tool's own code can work around.
    import fitz
    from openpyxl import load_workbook

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    x_cols = [72, 200, 400]
    rows = [
        ("Date", "Description", "Amount"),
        ("01/02", "Coffee Shop", "-4.50"),
        ("01/03", "Paycheck", "2500.00"),
    ]
    for i, row in enumerate(rows):
        for x, cell in zip(x_cols, row, strict=False):
            page.insert_text((x, 120 + i * 22), cell, fontsize=11)
    content = doc.tobytes()
    doc.close()

    xlsx_bytes = converter._convert_pdf_to_xlsx_sync(content)
    wb = load_workbook(io.BytesIO(xlsx_bytes))
    ws = wb["Page 1"]
    values = [
        tuple(r)
        for r in ws.iter_rows(values_only=True)
        if any(c not in (None, "") for c in r)
    ]
    # Real column separation, not one merged "Date Description Amount" cell.
    # The amount is a real negative float, not the text "-4.50" — proves the
    # formula-injection sanitizer (which only ever triggers on a leading
    # "=") didn't also mangle an ordinary negative number along the way.
    assert ("Date", "Description", "Amount") in values
    assert ("01/02", "Coffee Shop", -4.50) in values


def test_pdf_to_xlsx_neutralizes_formula_injection():
    # A crafted PDF whose "table" cell text starts with "=" must not become a
    # live Excel formula in the output — openpyxl auto-marks any raw
    # "="-prefixed string as data_type "f" (a real <f> formula element), so an
    # unsanitized cell here would execute in the victim's Excel on open
    # (CWE-1236). Goes through the plain-text fallback path since building a
    # ruled-table fixture with formula-looking cell text isn't necessary to
    # prove the sanitizer runs on every write path.
    import fitz
    from openpyxl import load_workbook

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 72), '=HYPERLINK("http://evil.example","x")', fontsize=12)
    content = doc.tobytes()
    doc.close()

    xlsx_bytes = converter._convert_pdf_to_xlsx_sync(content)
    wb = load_workbook(io.BytesIO(xlsx_bytes))
    ws = wb["Page 1"]
    cell = ws["A1"]
    assert cell.data_type == "s"  # plain string, not "f" (formula)
    assert cell.value.startswith("'=")


def test_sanitize_xlsx_cell_only_neutralizes_leading_equals():
    # Narrower than the usual CSV-injection "=+-@" prefix set on purpose —
    # openpyxl only auto-formula-types a leading "=" (verified directly
    # against Cell._bind_value), and "-" is the single most common leading
    # character in real extracted table data (negative amounts), so
    # sanitizing it too would corrupt legitimate numbers.
    assert converter._sanitize_xlsx_cell("=cmd|'/c calc'!A1") == "'=cmd|'/c calc'!A1"
    assert converter._sanitize_xlsx_cell("-4.50") == "-4.50"
    assert converter._sanitize_xlsx_cell("+1 555-1234") == "+1 555-1234"
    assert converter._sanitize_xlsx_cell("@handle") == "@handle"
    assert converter._sanitize_xlsx_cell("Widget") == "Widget"
    assert converter._sanitize_xlsx_cell(None) is None


def test_coerce_xlsx_numeric_converts_real_numbers_only():
    assert converter._coerce_xlsx_numeric("-4.50") == -4.50
    assert converter._coerce_xlsx_numeric("1,234") == 1234
    assert converter._coerce_xlsx_numeric("10") == 10
    assert isinstance(converter._coerce_xlsx_numeric("10"), int)
    assert isinstance(converter._coerce_xlsx_numeric("10.0"), float)
    # Not plausibly numeric — left as the original string.
    assert converter._coerce_xlsx_numeric("$5.00") == "$5.00"
    assert converter._coerce_xlsx_numeric("01/02") == "01/02"
    assert converter._coerce_xlsx_numeric("Widget") == "Widget"


def test_prepare_xlsx_cell_value_full_pipeline():
    # Numbers convert (and therefore never reach the string sanitizer at
    # all); non-numeric text only gets the leading-apostrophe treatment when
    # it starts with "=".
    assert converter._prepare_xlsx_cell_value("-4.50") == -4.50
    assert converter._prepare_xlsx_cell_value("=1+1") == "'=1+1"
    assert converter._prepare_xlsx_cell_value("Widget") == "Widget"
    assert converter._prepare_xlsx_cell_value(None) is None


def test_pdf_to_pptx_one_slide_per_page_matching_page_aspect_ratio():
    import fitz
    from pptx import Presentation

    doc = fitz.open()
    doc.new_page(width=595, height=842)  # portrait, A4-ish
    doc.new_page(width=595, height=842)
    content = doc.tobytes()
    doc.close()

    pptx_bytes = converter._convert_pdf_to_pptx_sync(content)
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = list(prs.slides)
    assert len(slides) == 2
    for slide in slides:
        shapes = list(slide.shapes)
        assert len(shapes) == 1
        assert shapes[0].shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
    # Deck sized to the source page's own aspect ratio, not python-pptx's
    # 10x7.5in default.
    assert abs(prs.slide_width / prs.slide_height - 595 / 842) < 0.01


def test_pdf_to_pptx_letterboxes_mismatched_aspect_ratio_page():
    import fitz
    from pptx import Presentation

    doc = fitz.open()
    doc.new_page(width=595, height=842)  # portrait — drives deck size
    doc.new_page(width=842, height=595)  # landscape page inside that deck
    content = doc.tobytes()
    doc.close()

    pptx_bytes = converter._convert_pdf_to_pptx_sync(content)
    prs = Presentation(io.BytesIO(pptx_bytes))
    landscape_slide = list(prs.slides)[1]
    pic = list(landscape_slide.shapes)[0]
    # Full width, vertically centered — not stretched to fill the whole slide.
    assert pic.width == prs.slide_width
    assert pic.height < prs.slide_height
    assert pic.left == 0
    assert abs(pic.top - (prs.slide_height - pic.height) / 2) <= 1


def _make_ring_png_bytes():
    from PIL import Image, ImageDraw

    img = Image.new("RGBA", (200, 200), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.ellipse((20, 20, 180, 180), fill=(220, 30, 30, 255))
    d.ellipse((70, 70, 130, 130), fill=(0, 0, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def test_trace_png_to_svg_ring_shape_has_correct_topology_and_color():
    svg_bytes = converter._trace_png_to_svg_sync(_make_ring_png_bytes())
    svg_text = svg_bytes.decode("utf-8")
    assert svg_text.count("<path") == 3
    assert 'fill="#dc1e1e"' in svg_text  # rgb(220,30,30), the ring's real color
    assert 'fill-rule="evenodd"' in svg_text
    assert 'viewBox="0 0 200 200"' in svg_text


def test_trace_png_to_svg_downscales_oversized_images():
    from PIL import Image

    big = Image.new("RGB", (3000, 1000), (10, 200, 10))
    buf = io.BytesIO()
    big.save(buf, format="PNG")

    svg_bytes = converter._trace_png_to_svg_sync(buf.getvalue())
    svg_text = svg_bytes.decode("utf-8")
    assert 'width="1500" height="500"' in svg_text


async def test_metrics_endpoint_shape(admin_client):
    r = await admin_client.get("/api/v1/metrics")
    assert r.status_code == 200
    body = r.json()
    assert set(body) == {"conversions", "failures", "avg_duration_ms"}


async def test_metrics_endpoint_requires_admin(client, user_client):
    # Anonymous: no session at all.
    assert (await client.get("/api/v1/metrics")).status_code == 401
    # Signed in, but not an admin.
    assert (await user_client.get("/api/v1/metrics")).status_code == 403


async def test_health_endpoint_responds(client):
    # Gotenberg host isn't resolvable from the test process → degraded, but the
    # endpoint must still respond 200 with a status.
    r = await client.get("/api/v1/health")
    assert r.status_code == 200
    body = r.json()
    assert body["status"] in ("healthy", "degraded")
    # P2 §20 — the DB check runs against the same test DB conftest wires every
    # other test to, so it's genuinely reachable here.
    assert body["database"] == "up"


async def test_health_endpoint_reports_db_down(client, monkeypatch):
    import converter

    class _BoomEngine:
        def connect(self):
            raise RuntimeError("connection refused")

    monkeypatch.setattr(converter, "async_engine", _BoomEngine())
    r = await client.get("/api/v1/health")
    assert r.status_code == 200
    body = r.json()
    assert body["database"] == "down"
    assert body["status"] == "degraded"


async def test_health_endpoint_bounds_a_hanging_db_connect(client, monkeypatch):
    # Post-merge audit fix. `async_engine` sets no connect_timeout (only the
    # sync engine does — data/db.py), so an unbounded `.connect()` could hang
    # past an uptime monitor's polling interval on a network partition that
    # drops packets rather than refusing outright — and unlike the one-time
    # startup check this mirrors, /health is now re-triggered by every poll,
    # so a hang here accumulates a stuck task per poll for as long as the
    # partition lasts. Proves the bound actually applies: a connect that
    # never completes still returns within the (shrunk, for a fast test)
    # timeout, degraded rather than hung.
    import asyncio

    import converter

    class _HangingConnectCM:
        async def __aenter__(self):
            await asyncio.sleep(10)  # would hang the request without the bound
            return self

        async def __aexit__(self, *exc):
            return False

    class _HangingEngine:
        def connect(self):
            return _HangingConnectCM()

    monkeypatch.setattr(converter, "async_engine", _HangingEngine())
    monkeypatch.setattr(converter, "HEALTH_DB_TIMEOUT_SECONDS", 0.05)

    # wait_for is the assertion, not just a safety net: if the internal bound
    # didn't apply, the hanging connect (sleep(10)) would blow past this and
    # fail loudly with TimeoutError instead of silently passing. 8s, not
    # ~0.05s, because the Gotenberg check runs concurrently (asyncio.gather)
    # and unrelatedly — its httpx timeout only bounds the CONNECT phase, and a
    # DNS failure for the unresolvable test-env hostname can itself take a
    # couple of seconds; the margin absorbs that without weakening what this
    # test actually proves (nowhere near the 10s the hanging connect would
    # take unbounded).
    r = await asyncio.wait_for(client.get("/api/v1/health"), timeout=8.0)
    assert r.status_code == 200
    body = r.json()
    assert body["database"] == "down"
    assert body["status"] == "degraded"


async def test_health_endpoint_accepts_head(client):
    # Uptime monitors default to HEAD requests; must not 405.
    r = await client.head("/api/v1/health")
    assert r.status_code == 200
    assert r.content == b""


# --------------------------------------------------------------------------- #
# Gotenberg request queue (P3 §31) — caps concurrent Gotenberg calls with an
# in-process asyncio.Semaphore. `_gotenberg_request` is exercised directly
# (not through the route) with a fake httpx.AsyncClient, since these tests
# assert on timing/concurrency, not on any particular tool's request shape.
# --------------------------------------------------------------------------- #


class _FakeResponse:
    def __init__(self, status_code=200, content=b"ok"):
        self.status_code = status_code
        self.content = content
        self.text = content.decode()


def _fake_async_client(post_impl):
    class _FakeAsyncClient:
        def __init__(self, *a, **kw):
            pass

        async def __aenter__(self):
            return self

        async def __aexit__(self, *exc):
            return False

        async def post(self, url, files=None, headers=None):
            return await post_impl()

    return _FakeAsyncClient


async def test_gotenberg_semaphore_caps_concurrency(monkeypatch):
    monkeypatch.setattr(converter, "_gotenberg_semaphore", asyncio.Semaphore(1))
    monkeypatch.setattr(converter, "GOTENBERG_QUEUE_TIMEOUT_SECONDS", 5.0)

    in_flight = 0
    max_in_flight = 0

    async def post_impl():
        nonlocal in_flight, max_in_flight
        in_flight += 1
        max_in_flight = max(max_in_flight, in_flight)
        await asyncio.sleep(0.05)
        in_flight -= 1
        return _FakeResponse()

    monkeypatch.setattr(converter.httpx, "AsyncClient", _fake_async_client(post_impl))

    # Three "requests" sharing a single-slot semaphore must never overlap,
    # even though they're all launched at once.
    await asyncio.gather(
        *(
            converter._gotenberg_request("/forms/libreoffice/convert", {}, "test")
            for _ in range(3)
        )
    )
    assert max_in_flight == 1


async def test_gotenberg_queue_timeout_returns_503_with_retry_after(
    client, monkeypatch
):
    monkeypatch.setattr(converter, "_gotenberg_semaphore", asyncio.Semaphore(1))
    monkeypatch.setattr(converter, "GOTENBERG_QUEUE_TIMEOUT_SECONDS", 0.05)

    async def post_impl():
        # Holds the one slot well past the queue timeout so the second
        # request below is forced to wait and expire.
        await asyncio.sleep(1.0)
        return _FakeResponse()

    monkeypatch.setattr(converter.httpx, "AsyncClient", _fake_async_client(post_impl))

    valid_docx = b"PK\x03\x04" + b"\x00" * 200

    async def _convert_request():
        return await client.post(
            "/api/v1/convert/docx-to-pdf",
            files={
                "file": (
                    "report.docx",
                    valid_docx,
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
        )

    # Fire two conversions at once: the first occupies the only slot for 1s
    # (well past the 0.05s queue timeout), so the second must time out
    # waiting rather than queue indefinitely.
    first, second = await asyncio.gather(_convert_request(), _convert_request())
    timed_out = [r for r in (first, second) if r.status_code == 503]
    assert timed_out, (first.status_code, second.status_code)
    body = timed_out[0].json()
    assert body["error_type"] == "queue_timeout"
    assert timed_out[0].headers["retry-after"] == "10"
